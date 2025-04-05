import os
import requests
import json
import docx
import pptx
import PyPDF2
from flask import Flask, request, render_template, redirect, url_for, session, jsonify

app = Flask(__name__)
app.secret_key = os.environ.get('SECRET_KEY', 'supersecretkey')  # Get secret key from env

# GROQ_API_KEY = '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' # Removed hardcoded key
MODEL_NAME = 'llama3-70b-8192'

ALLOWED_EXTENSIONS = {'txt', 'pdf', 'doc', 'docx', 'ppt', 'pptx'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_file(uploaded_file, ext):
    if ext == 'txt':
        return uploaded_file.read().decode('utf-8')
    elif ext == 'pdf':
        reader = PyPDF2.PdfReader(uploaded_file)
        return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
    elif ext in ['doc', 'docx']:
        doc = docx.Document(uploaded_file)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    elif ext in ['ppt', 'pptx']:
        prs = pptx.Presentation(uploaded_file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    else:
        raise ValueError("Unsupported file extension")

def generate_quiz_from_text(text, api_key, model):
    prompt = f"""
    Generate 30 multiple-choice questions from the following content.
    Each question should have exactly 4 options with one correct answer marked.
    Format:
    Q1. <question text>
    A. <option>
    B. <option>
    C. <option>
    D. <option>
    Answer: B
    Content: '''{text}'''
    """

    response = requests.post(
        "https://api.groq.com/openai/v1/chat/completions",
        headers={"Authorization": f"Bearer {api_key}"},
        json={
            "model": model,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.7
        }
    )

    result_text = response.json()['choices'][0]['message']['content']
    return parse_questions(result_text)

def parse_questions(raw_text):
    questions = []
    blocks = raw_text.strip().split("Q")[1:]
    for block in blocks:
        lines = block.strip().split("\n")
        if len(lines) < 6:
            continue  # Skip malformed blocks
        q_text = lines[0].strip().replace("**", "")  # Remove ** from questions
        options = [line.strip()[3:] for line in lines[1:5] if len(line.strip()) > 2]
        answer_lines = [line for line in lines if line.strip().lower().startswith("answer")]
        if not answer_lines or len(options) != 4:
            continue  # Skip if format is wrong
        correct = answer_lines[0].strip()[-1].upper()
        questions.append({
            "question": q_text,
            "options": options,
            "answer": correct
        })
    return questions

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload():
    uploaded_file = request.files['file']
    if uploaded_file and allowed_file(uploaded_file.filename):
        ext = uploaded_file.filename.rsplit('.', 1)[1].lower()

        try:
            text_content = extract_text_from_file(uploaded_file, ext)
        except Exception as e:
            return f"Failed to extract text from file: {str(e)}"

        groq_api_key = os.environ.get('GROQ_API_KEY')
        if not groq_api_key:
            return "Error: GROQ_API_KEY environment variable not set."

        questions = generate_quiz_from_text(text_content, groq_api_key, MODEL_NAME)

        os.makedirs('static', exist_ok=True)
        with open('static/quiz_data.json', 'w', encoding='utf-8') as f:
            json.dump(questions, f, indent=2)

        return redirect(url_for('quiz')) # Corrected redirect

    return "Unsupported file type or no file uploaded."

# @app.route('/quiz')
# def quiz():
#     return redirect(url_for('quiz')) # Redirect /quiz to /quezz

@app.route('/quiz')
def quiz():
    return render_template('quiz.html') # Corrected template name

# 🔥 Store results from frontend
@app.route('/leaderboard', methods=['POST'])
def save_results():
    data = request.json
    session['result'] = data
    return jsonify({'message': 'Result saved successfully'})

# 🔥 Serve result for leaderboard.html
@app.route('/get_results')
def get_results():
    result = session.get('result', {
        'totalScore': 0,
        'correctAnswers': 0,
        'timeTaken': '0s'
    })
    return jsonify(result)

# 👇 If you want to serve leaderboard.html as a GET (optional)
@app.route('/leaderboard.html')
def leaderboard_html():
    return render_template('leaderboard.html')

@app.route('/update_leaderboard', methods=['POST'])
def update_leaderboard():
    data = request.json
    # Process and store data
    return jsonify({"message": "Leaderboard updated successfully"})

# Removed unwanted route
# @app.route('/hybridaction/zybTrackerStatisticsAction', methods=['GET', 'POST'])
# def block_unwanted_requests(subpath):
#     return jsonify({"error": "Invalid request"}), 404

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=int(os.environ.get('PORT', 5000))) # Listen on all interfaces and dynamic port
